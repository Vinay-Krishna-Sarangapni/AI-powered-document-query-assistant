import os
import time
import io
import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from PyPDF2 import PdfReader
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain.vectorstores import FAISS
from dotenv import load_dotenv
import google.generativeai as genai

# Load environment variables from .env file
load_dotenv()
google_api_key = os.getenv("GOOGLE_GEN_AI_KEY")

# Configure Google Generative AI
genai.configure(api_key=google_api_key)

def get_pdf_text(pdf_file):
    """Extract text from a PDF file."""
    text = ""
    try:
        pdf_reader = PdfReader(pdf_file)
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        st.error(f"Error processing PDF {pdf_file.name}: {str(e)}")
    return text

def get_ppt_text(ppt_file):
    """Extract text from a PowerPoint file."""
    text = ""
    try:
        ppt_bytes = ppt_file.read()
        ppt_file_io = io.BytesIO(ppt_bytes)
        ppt_reader = Presentation(ppt_file_io)
        for slide in ppt_reader.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        st.error(f"Error processing PowerPoint {ppt_file.name}: {str(e)}")
    finally:
        ppt_file.seek(0)  # Reset file pointer
    return text

def get_text_chunks(text):
    """Split text into meaningful chunks."""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=2000,
        chunk_overlap=200,
        separators=["\n\n", "\n", ".", "?", "!"]
    )
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks):
    """Create and save a FAISS vector store."""
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def semantic_search(query, k=15):
    """Perform semantic search using FAISS."""
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    return db.similarity_search(query, k=k)

def improved_rag(query, context, answer_type, model_name):
    """Improved RAG method with better context utilization."""
    if answer_type == "Detailed":
        prompt = f"""
        Based on the following context, provide a comprehensive and detailed answer to the question. 
        Your response should be thorough, well-structured, and demonstrate a deep understanding of the subject matter.
        If the information is not sufficient to answer the question fully, provide the most relevant information available 
        and explain any limitations or uncertainties.

        Context: {context}

        Question: {query}

        Detailed Answer:
        """
    else:  # Brief
        prompt = f"""
        Based on the following context, provide a brief but informative summary answer to the question.
        Capture the key points and main ideas while maintaining clarity and coherence.
        If the information is not sufficient to answer the question, provide a concise explanation of what is known.

        Context: {context}

        Question: {query}

        Brief Answer:
        """

    # Select model based on user input
    model = ChatGoogleGenerativeAI(model=model_name, temperature=0.2)
    response = model.predict(prompt)
    
    return response

def search_and_answer(user_question, answer_type, model_name):
    """Simulate a search and answer function for a specific model."""
    time.sleep(1)  # Simulate processing time
    response = f"{model_name} response to '{user_question}' with {answer_type} answer."
    return response

def semantic_search_and_answer(user_question, answer_type, model_name):
    """Search and generate an answer using improved RAG."""
    try:
        relevant_docs = semantic_search(user_question, k=15)
        context = "\n".join([doc.page_content for doc in relevant_docs])
        
        answer = improved_rag(user_question, context, answer_type, model_name)
        
        st.write(f"{'Detailed' if answer_type == 'Detailed' else 'Brief'} Answer:")
        st.write(answer)

    except Exception as e:
        st.error(f"An error occurred while processing your question: {str(e)}")

def external_search(query, answer_type, model_name):
    """Perform an external search using Google Generative AI."""
    try:
        model = genai.GenerativeModel(model_name)
        if answer_type == "Detailed":
            prompt = f"""
            Provide a comprehensive and detailed answer to the following question:
            {query}
            Your response should be thorough, well-structured, and demonstrate a deep understanding of the subject matter.
            Include relevant examples, explanations, and any necessary context to fully address the question.
            """
        else:  # Brief
            prompt = f"""
            Provide a brief but informative summary answer to the following question:
            {query}
            Capture the key points and main ideas while maintaining clarity and coherence.
            """
        
        response = model.generate_content(prompt)
        
        if response.text:
            return response.text
        else:
            return "Unable to generate an answer."
    except Exception as e:
        st.error(f"Failed to perform search: {str(e)}")
        return "Error during search."

def user_input(user_question, search_type, answer_type, model_name):
    """Process user input based on selected search type and answer type."""
    if search_type == "Document Search":
        st.write(f"Generating {'detailed' if answer_type == 'Detailed' else 'brief'} answer using {search_type}...")
        semantic_search_and_answer(user_question, answer_type, model_name)
    
    elif search_type == "External Search":
        st.write(f"Performing external search for {'detailed' if answer_type == 'Detailed' else 'brief'} answer...")
        with st.spinner("Generating answer..."):
            search_result = external_search(user_question, answer_type, model_name)
            if search_result:  # Check for valid response
                st.write(f"{'Detailed' if answer_type == 'Detailed' else 'Brief'} Answer:")
                st.write(search_result)
            else:
                st.warning("No valid response returned from the external search.")


def process_files(uploaded_files):
    """Process uploaded files and create a vector store."""
    all_text = ""
    progress_bar = st.progress(0)
    for i, file in enumerate(uploaded_files):
        if file.type == "application/pdf":
            all_text += get_pdf_text(file)
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            all_text += get_ppt_text(file)
        else:
            st.warning(f"Unsupported file type: {file.type}")
        progress_bar.progress((i + 1) / len(uploaded_files))

    if all_text:
        text_chunks = get_text_chunks(all_text)
        get_vector_store(text_chunks)
        st.success("All files processed successfully!")
    else:
        st.warning("No text could be extracted from the uploaded files.")

def main():
    """Main function to run the Streamlit app."""
    st.set_page_config(page_title="Advanced Document Chat Assistant and Model Comparison Tool")
    
    # Sidebar for file upload and model comparison
    st.sidebar.header("Document Upload")
    uploaded_files = st.sidebar.file_uploader("Upload Document Files", type=["pdf", "pptx"], accept_multiple_files=True)
    if st.sidebar.button("Process Files"):
        if uploaded_files:
            process_files(uploaded_files)
        else:
            st.warning("Please upload at least one document.")
    
    st.header("LLM Chatbot for Comprehensive Document Assistance")
    
    # User Input for Chatbot
    search_type = st.selectbox("Select Search Type", ["Document Search", "External Search"])
    answer_type = st.selectbox("Select Answer Type", ["Brief", "Detailed"])
    model_name = st.selectbox("Select Model", ["gemini-pro", "gemini-1.5-pro"])

    user_question = st.text_input("Ask a Question")

    if st.button("Submit Question"):
        if user_question:
            user_input(user_question, search_type, answer_type, model_name)
        else:
            st.warning("Please enter a question.")

    # Model Comparison Tool
    st.header("Model Comparison Tool")
    user_question_compare = st.text_input("Enter your question for model comparison:")
    answer_type_compare = st.selectbox("Select answer type for comparison:", ["Detailed", "Brief"], key="compare")
    
    models = ["gemini-pro", "gemini-1.5-pro", "model-3"]
    selected_models = st.multiselect("Select two models to compare:", models, default=["gemini-pro", "gemini-1.5-pro"], max_selections=2)

    if st.button("Get Answers for Comparison"):
        if len(selected_models) == 2:
            st.write(f"Generating {answer_type_compare} answers for the selected models...")
            
            performance_data = []
            for model_name in selected_models:
                start_time = time.time()  # Start time for performance measurement
                
                # Fetch answer using selected model
                if search_type == "Document Search":
                    response = search_and_answer(user_question_compare, answer_type_compare, model_name)
                else:
                    response = external_search(user_question_compare, answer_type_compare, model_name)

                end_time = time.time()  # End time for performance measurement
                response_time = end_time - start_time  # Calculate response time
                response_length = len(response)  # Measure response length
                
                # Store performance metrics
                performance_data.append({
                    "Model": model_name,
                    "Response": response,
                    "Response Time (s)": response_time,
                    "Response Length": response_length
                })
                
                # Display the answer for each model
                st.subheader(f"{model_name} Answer:")
                st.write(response)
        
            # Convert performance data to a DataFrame for visualization
            df_performance = pd.DataFrame(performance_data)
            
            # Visualization
            st.subheader("Performance Metrics")
            st.write(df_performance)
            
            # Plotting the performance metrics
            fig, ax1 = plt.subplots()

            # Bar chart for response time
            color = 'tab:blue'
            ax1.set_xlabel('Models')
            ax1.set_ylabel('Response Time (s)', color=color)
            ax1.bar(df_performance['Model'], df_performance['Response Time (s)'], color=color, alpha=0.6, label='Response Time')
            ax1.tick_params(axis='y', labelcolor=color)

            # Creating a second y-axis for response length
            ax2 = ax1.twinx()
            color = 'tab:red'
            ax2.set_ylabel('Response Length', color=color)
            ax2.plot(df_performance['Model'], df_performance['Response Length'], color=color, marker='o', label='Response Length')
            ax2.tick_params(axis='y', labelcolor=color)

            # Adding a title and legend
            plt.title('Model Performance Comparison')
            fig.tight_layout()  # Adjust layout to prevent clipping
            st.pyplot(fig)
        else:
            st.warning("Please select exactly two models for comparison.")

if __name__ == "__main__":
    main()   
